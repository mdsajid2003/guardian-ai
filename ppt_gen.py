"""
ppt_gen.py — Guardian AI: PowerPoint report generator
4 slides: Title | Executive Summary | Leakage Map | AI Action Plan
Every slide includes a "View Source" note as per the architecture spec.
"""

import io
from datetime import datetime
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from state import AnalysisResult, AmazonReconResult

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x12, 0x12, 0x24)
CARD_BG      = RGBColor(0x1E, 0x1E, 0x38)
ACCENT_GREEN = RGBColor(0x00, 0xB4, 0x8A)
ACCENT_ORANGE= RGBColor(0xFF, 0x6B, 0x35)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
MUTED        = RGBColor(0x88, 0x88, 0xAA)


# ── Helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_filled_rect(slide, left, top, width, height, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color  # match border to fill = invisible border


def _source_footer(slide) -> None:
    """Every slide gets a 'View Source' link placeholder."""
    _add_text(
        slide,
        "View Source: Guardian AI · guardian-ai.streamlit.app",
        0.3, 7.1, 12.7, 0.3,
        size=9, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ── Slide builders ────────────────────────────────────────────────────────────

def _slide_title(prs: Presentation, company: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)

    # Accent bar
    _add_filled_rect(slide, 0, 3.1, 13.33, 0.06, ACCENT_GREEN)

    _add_text(slide, "Guardian AI", 0, 0.9, 13.33, 1.4,
              size=52, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    _add_text(slide, "Agentic Accountant  ·  Financial Intelligence Report",
              0, 2.4, 13.33, 0.7,
              size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    _add_text(slide, company,
              0, 3.3, 13.33, 0.6,
              size=17, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide, f"Generated {datetime.now().strftime('%d %B %Y')}",
              0, 4.1, 13.33, 0.5,
              size=13, color=MUTED, align=PP_ALIGN.CENTER)
    _source_footer(slide)


def _slide_exec_summary(prs: Presentation, analysis: AnalysisResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_filled_rect(slide, 0.4, 0.25, 12.5, 0.06, ACCENT_GREEN)

    _add_text(slide, "02  Executive Summary",
              0.5, 0.45, 9, 0.65, size=26, bold=True, color=ACCENT_GREEN)
    _add_text(slide, "Slide 2 / 4", 10.5, 0.5, 2.3, 0.5,
              size=11, color=MUTED, align=PP_ALIGN.RIGHT)

    # Summary box
    _add_filled_rect(slide, 0.5, 1.25, 12.3, 0.75, CARD_BG)
    _add_text(slide, f"📊  {analysis.summary}",
              0.7, 1.32, 11.9, 0.65, size=16, bold=True, color=WHITE)

    # Observations
    _add_text(slide, "KEY OBSERVATIONS", 0.5, 2.2, 6, 0.4,
              size=11, bold=True, color=ACCENT_ORANGE)
    obs_lines = "\n".join(f"▸  {o}" for o in (analysis.observations or [])[:5])
    _add_text(slide, obs_lines or "—", 0.5, 2.65, 6, 2.8,
              size=13, color=LIGHT_GRAY)

    # Focused question callout
    if analysis.focused_question:
        _add_filled_rect(slide, 7.1, 2.2, 5.7, 1.0, CARD_BG)
        _add_text(slide, f"💡  {analysis.focused_question}",
                  7.2, 2.28, 5.5, 0.9, size=13, color=ACCENT_GREEN)

    _source_footer(slide)


def _slide_leakage(prs: Presentation, recon: Optional[AmazonReconResult]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_filled_rect(slide, 0.4, 0.25, 12.5, 0.06, ACCENT_ORANGE)

    _add_text(slide, "03  Revenue Leakage Map",
              0.5, 0.45, 9, 0.65, size=26, bold=True, color=ACCENT_ORANGE)
    _add_text(slide, "Slide 3 / 4", 10.5, 0.5, 2.3, 0.5,
              size=11, color=MUTED, align=PP_ALIGN.RIGHT)

    if recon and (recon.mtr_total or recon.leakage_amount):
        # Metric cards row
        metrics = [
            ("MTR Total",     f"₹{recon.mtr_total:,.0f}"       if recon.mtr_total        else "N/A"),
            ("Settlement",    f"₹{recon.settlement_total:,.0f}" if recon.settlement_total else "N/A"),
            ("Leakage",       f"₹{recon.leakage_amount:,.0f}"   if recon.leakage_amount   else "N/A"),
            ("ACOS",          f"{recon.acos:.1f}%"              if recon.acos             else "N/A"),
        ]
        for i, (label, value) in enumerate(metrics):
            x = 0.5 + i * 3.1
            _add_filled_rect(slide, x, 1.35, 2.8, 1.1, CARD_BG)
            _add_text(slide, label, x + 0.1, 1.42, 2.6, 0.4,
                      size=11, color=MUTED)
            _add_text(slide, value, x + 0.1, 1.78, 2.6, 0.55,
                      size=20, bold=True, color=WHITE)

        # GST breakdown
        gst = recon.gst_reconciliation
        if gst:
            _add_text(slide, "GST RECONCILIATION", 0.5, 2.65, 5.8, 0.35,
                      size=11, bold=True, color=ACCENT_GREEN)
            gst_lines = "\n".join([
                f"IGST Collected:    ₹{gst.get('igst_collected', 0):,.2f}",
                f"CGST Collected:    ₹{gst.get('cgst_collected', 0):,.2f}",
                f"SGST Collected:    ₹{gst.get('sgst_collected', 0):,.2f}",
                f"GST in Settlement: ₹{gst.get('gst_in_settlement', 0):,.2f}",
                f"GST Gap:           ₹{gst.get('gst_gap', 0):,.2f}",
            ])
            _add_text(slide, gst_lines, 0.5, 3.05, 5.8, 2.5,
                      size=12, color=LIGHT_GRAY)

        # Fee breakdown
        fees = recon.fee_breakdown
        if fees:
            _add_text(slide, "FBA FEE BREAKDOWN", 7.1, 2.65, 5.7, 0.35,
                      size=11, bold=True, color=ACCENT_ORANGE)
            fee_lines = "\n".join(
                f"{k.replace('_', ' ').title()}: ₹{v:,.2f}" for k, v in fees.items()
            )
            _add_text(slide, fee_lines, 7.1, 3.05, 5.7, 2.5,
                      size=12, color=LIGHT_GRAY)

        # Recommendations
        if recon.recommendations:
            recs = "\n".join(f"✓  {r}" for r in recon.recommendations[:3])
            _add_filled_rect(slide, 0.5, 5.7, 12.3, 1.1, CARD_BG)
            _add_text(slide, recs, 0.7, 5.78, 11.9, 0.95,
                      size=12, color=ACCENT_GREEN)
    else:
        _add_text(
            slide,
            "Upload MTR and Settlement reports in the Amazon Reconciliation tab\n"
            "to enable leakage detection and GST gap analysis.",
            0.5, 2.5, 12.3, 2, size=15, color=MUTED, align=PP_ALIGN.CENTER,
        )

    _source_footer(slide)


def _slide_action_plan(prs: Presentation, analysis: AnalysisResult) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)
    _add_filled_rect(slide, 0.4, 0.25, 12.5, 0.06, ACCENT_GREEN)

    _add_text(slide, "04  AI Action Plan",
              0.5, 0.45, 9, 0.65, size=26, bold=True, color=ACCENT_GREEN)
    _add_text(slide, "Slide 4 / 4", 10.5, 0.5, 2.3, 0.5,
              size=11, color=MUTED, align=PP_ALIGN.RIGHT)

    # Insights column
    _add_text(slide, "INSIGHTS", 0.5, 1.3, 6, 0.4,
              size=11, bold=True, color=ACCENT_ORANGE)
    ins_lines = "\n".join(f"▸  {ins}" for ins in (analysis.insights or [])[:5])
    _add_text(slide, ins_lines or "—", 0.5, 1.75, 6, 3.5,
              size=13, color=LIGHT_GRAY)

    # Action items column
    _add_text(slide, "ACTION ITEMS", 7.1, 1.3, 5.7, 0.4,
              size=11, bold=True, color=ACCENT_GREEN)
    act_lines = "\n".join(f"☐  {act}" for act in (analysis.action_items or [])[:5])
    _add_text(slide, act_lines or "—", 7.1, 1.75, 5.7, 3.5,
              size=13, color=LIGHT_GRAY)

    # Bottom callout
    if analysis.focused_question:
        _add_filled_rect(slide, 0.5, 5.5, 12.3, 1.1, CARD_BG)
        _add_text(slide, f"Next Question to Resolve:  💡  {analysis.focused_question}",
                  0.7, 5.6, 11.9, 0.9, size=14, bold=True,
                  color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    _source_footer(slide)


# ── Public entry point ────────────────────────────────────────────────────────

def create_presentation(
    analysis: AnalysisResult,
    amazon_recon: Optional[AmazonReconResult] = None,
    company_name: str = "Guardian AI Report",
) -> bytes:
    """
    Build a 4-slide PowerPoint and return raw bytes.
    Slide 1 – Title
    Slide 2 – Executive Summary
    Slide 3 – Revenue Leakage Map
    Slide 4 – AI Action Plan
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _slide_title(prs, company_name)
    _slide_exec_summary(prs, analysis)
    _slide_leakage(prs, amazon_recon)
    _slide_action_plan(prs, analysis)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
